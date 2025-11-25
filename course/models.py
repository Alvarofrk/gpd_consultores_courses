from django.conf import settings
from django.utils.text import slugify
from django.core.validators import FileExtensionValidator
from django.db import models
from django.db.models import Q
from django.db.models.signals import pre_save, post_delete, post_save
from django.dispatch import receiver
from django.urls import reverse
from django.utils.translation import gettext_lazy as _

from core.models import ActivityLog, Semester
from core.utils import unique_slug_generator
from django.contrib.auth.models import User
import re  # Importamos el módulo 're' para usar expresiones regulares
import os
import zipfile
import io
import base64
from django.core.exceptions import ValidationError
from pptx import Presentation


class ProgramManager(models.Manager):
    def search(self, query=None):
        queryset = self.get_queryset()
        if query:
            or_lookup = Q(title__icontains=query) | Q(summary__icontains=query)
            queryset = queryset.filter(or_lookup).distinct()
        return queryset


class Program(models.Model):
    title = models.CharField(max_length=150, unique=True)
    summary = models.TextField(blank=True)

    objects = ProgramManager()

    def __str__(self):
        return f"{self.title}"

    def get_absolute_url(self):
        return reverse("program_detail", kwargs={"pk": self.pk})


@receiver(post_save, sender=Program)
def log_program_save(sender, instance, created, **kwargs):
    verb = "created" if created else "updated"
    ActivityLog.objects.create(message=_(f"The program '{instance}' has been {verb}."))


@receiver(post_delete, sender=Program)
def log_program_delete(sender, instance, **kwargs):
    ActivityLog.objects.create(message=_(f"The program '{instance}' has been deleted."))


class CourseManager(models.Manager):
    def search(self, query=None):
        queryset = self.get_queryset()
        if query:
            or_lookup = (
                Q(title__icontains=query)
                | Q(summary__icontains=query)
                | Q(code__icontains=query)
                | Q(slug__icontains=query)
            )
            queryset = queryset.filter(or_lookup).distinct()
        return queryset



class Course(models.Model):
    slug = models.SlugField(unique=True, blank=True, max_length=200)
    title = models.CharField(max_length=200)
    code = models.CharField(max_length=200, unique=True)
    credit = models.IntegerField(default=0)
    summary = models.TextField(max_length=200, blank=True)
    program = models.ForeignKey(Program, on_delete=models.CASCADE)
    level = models.CharField(max_length=25, choices=settings.LEVEL_CHOICES)
    year = models.IntegerField(choices=settings.YEARS, default=1)
    semester = models.CharField(choices=settings.SEMESTER_CHOICES, max_length=200)
    is_elective = models.BooleanField(default=False)
    is_active = models.BooleanField(default=True)
    is_external = models.BooleanField(
        default=False,
        verbose_name=_("Curso Externo"),
        help_text=_("Indica si este es un curso externo (certificado en Google Drive)")
    )

    last_cert_code = models.IntegerField(default=0)

    objects = CourseManager()

    def __str__(self):
        return f"{self.title} ({self.code})"

    def get_absolute_url(self):
        return reverse("course_detail", kwargs={"slug": self.slug})

    def save(self, *args, **kwargs):
        if not self.slug:
            self.slug = slugify(self.title)
        super().save(*args, **kwargs)

    @property
    def is_current_semester(self):
        current_semester = Semester.objects.filter(is_current_semester=True).first()
        return self.semester == current_semester.semester if current_semester else False

    def get_progress_for_user(self, user):
        """Calcula el progreso del curso para un usuario específico"""
        videos = UploadVideo.objects.filter(course=self)
        documents = Upload.objects.filter(course=self)
        
        total_content = videos.count() + documents.count()
        if total_content == 0:
            return 0
        
        completed_videos = VideoCompletion.objects.filter(
            user=user, video__in=videos
        ).count()
        completed_docs = DocumentCompletion.objects.filter(
            user=user, document__in=documents
        ).count()
        
        completed_content = completed_videos + completed_docs
        return round((completed_content / total_content) * 100, 1)
    
    def get_content_summary(self):
        """Obtiene resumen del contenido del curso"""
        videos = UploadVideo.objects.filter(course=self)
        documents = Upload.objects.filter(course=self)
        
        return {
            'total_videos': videos.count(),
            'total_documents': documents.count(),
            'total_content': videos.count() + documents.count(),
            'has_videos': videos.exists(),
            'has_documents': documents.exists(),
        }
    
    def get_user_completion_summary(self, user):
        """Obtiene resumen de completado para un usuario"""
        videos = UploadVideo.objects.filter(course=self)
        documents = Upload.objects.filter(course=self)
        
        completed_videos = VideoCompletion.objects.filter(
            user=user, video__in=videos
        ).count()
        completed_docs = DocumentCompletion.objects.filter(
            user=user, document__in=documents
        ).count()
        
        return {
            'completed_videos': completed_videos,
            'completed_documents': completed_docs,
            'completed_content': completed_videos + completed_docs,
            'total_videos': videos.count(),
            'total_documents': documents.count(),
            'total_content': videos.count() + documents.count(),
        }
    
    def get_course_status_for_user(self, user):
        """Determina el estado completo del curso para un usuario"""
        # 1. Verificar progreso del material
        material_progress = self.get_progress_for_user(user)
        
        # 2. Verificar si hay examen para este curso
        from quiz.models import Quiz, Sitting
        quiz = Quiz.objects.filter(course=self, draft=False).first()
        
        if not quiz:
            # Sin examen, solo material
            if material_progress == 100:
                return 'course_completed'
            elif material_progress > 0:
                return 'material_in_progress'
            else:
                return 'not_started'
        
        # 3. Con examen - verificar estado
        if material_progress < 100:
            return 'material_in_progress'
        
        # Material completado, verificar examen
        approved_sitting = Sitting.objects.filter(
            user=user,
            quiz=quiz,
            course=self,
            complete=True,
            current_score__gte=quiz.pass_mark * quiz.get_max_score / 100
        ).first()
        
        if approved_sitting:
            return 'course_completed'  # Material + Examen aprobado
        
        # Verificar si intentó examen
        attempted_sitting = Sitting.objects.filter(
            user=user,
            quiz=quiz,
            course=self,
            complete=True
        ).exists()
        
        if attempted_sitting:
            return 'exam_failed'  # Intentó pero no aprobó
        else:
            return 'exam_available'  # Puede hacer examen
    
    def get_exam_info_for_user(self, user):
        """Obtiene información detallada del examen para un usuario"""
        from quiz.models import Quiz, Sitting
        quiz = Quiz.objects.filter(course=self, draft=False).first()
        
        if not quiz:
            return None
        
        # Buscar intentos del usuario
        sittings = Sitting.objects.filter(
            user=user,
            quiz=quiz,
            course=self,
            complete=True
        ).order_by('-end')
        
        if not sittings.exists():
            return {
                'has_exam': True,
                'attempted': False,
                'passed': False,
                'best_score': 0,
                'attempts': 0,
                'can_retake': True,
                'pass_mark': quiz.pass_mark,
                'quiz_title': quiz.title
            }
        
        best_sitting = sittings.first()
        passed = best_sitting.check_if_passed
        
        return {
            'has_exam': True,
            'attempted': True,
            'passed': passed,
            'best_score': best_sitting.get_percent_correct,
            'attempts': sittings.count(),
            'can_retake': not quiz.single_attempt and not passed,
            'pass_mark': quiz.pass_mark,
            'quiz_title': quiz.title,
            'last_attempt': best_sitting.end,
            'fecha_aprobacion': best_sitting.fecha_aprobacion if passed else None
        }


@receiver(pre_save, sender=Course)
def course_pre_save_receiver(sender, instance, **kwargs):
    if not instance.slug:
        instance.slug = unique_slug_generator(instance)


@receiver(post_save, sender=Course)
def log_course_save(sender, instance, created, **kwargs):
    verb = "created" if created else "updated"
    ActivityLog.objects.create(message=_(f"The course '{instance}' has been {verb}."))


@receiver(post_delete, sender=Course)
def log_course_delete(sender, instance, **kwargs):
    ActivityLog.objects.create(message=_(f"The course '{instance}' has been deleted."))


class CourseAllocation(models.Model):
    lecturer = models.ForeignKey(
        settings.AUTH_USER_MODEL,
        on_delete=models.CASCADE,
        related_name="allocated_lecturer",
    )
    courses = models.ManyToManyField(Course, related_name="allocated_course")
    session = models.ForeignKey(
        "core.Session", on_delete=models.CASCADE, blank=True, null=True
    )
 
    def __str__(self):
        return self.lecturer.get_full_name

    def get_absolute_url(self):
        return reverse("edit_allocated_course", kwargs={"pk": self.pk})


class Upload(models.Model):
    title = models.CharField(max_length=100)
    course = models.ForeignKey(Course, on_delete=models.CASCADE)
    file = models.FileField(
        upload_to="course_files/",
        help_text=_(
            "Valid Files: pdf, docx, doc, xls, xlsx, ppt, pptx, zip, rar, 7zip"
        ),
        validators=[
            FileExtensionValidator(
                [
                    "pdf",
                    "docx",
                    "doc",
                    "xls",
                    "xlsx",
                    "ppt",
                    "pptx",
                    "zip",
                    "rar",
                    "7zip",
                ]
            )
        ],
        blank=True,
        null=True,
    )
    external_url = models.URLField(
        max_length=500,
        blank=True,
        null=True,
        help_text=_("URL externa del archivo (ej: Google Drive, Dropbox, etc.)")
    )
    updated_date = models.DateTimeField(auto_now=True)
    upload_time = models.DateTimeField(auto_now_add=True)

    def __str__(self):
        return f"{self.title}"

    def get_extension_short(self):
        if self.external_url:
            # Para URLs externas, intentar detectar el tipo por la URL
            if 'drive.google.com' in self.external_url:
                # Para Google Drive, asumir PDF por defecto ya que es el más común
                # y se puede visualizar en iframe
                return "pdf"
            elif 'dropbox.com' in self.external_url:
                # Para Dropbox, detectar por extensión en la URL
                if '.pdf' in self.external_url.lower():
                    return "pdf"
                elif '.ppt' in self.external_url.lower() or '.pptx' in self.external_url.lower():
                    return "powerpoint"
                else:
                    return "file"
            else:
                # Para otras URLs, detectar por extensión
                if '.pdf' in self.external_url.lower():
                    return "pdf"
                elif '.ppt' in self.external_url.lower() or '.pptx' in self.external_url.lower():
                    return "powerpoint"
                else:
                    return "file"
            return "file"
        
        if self.file and hasattr(self.file, 'name') and self.file.name:
            # Obtener la extensión del archivo
            file_name = self.file.name.lower()
            
            # Detectar por extensión
            if file_name.endswith(('.doc', '.docx')):
                return "word"
            elif file_name.endswith('.pdf'):
                return "pdf"
            elif file_name.endswith(('.xls', '.xlsx')):
                return "excel"
            elif file_name.endswith(('.ppt', '.pptx')):
                return "powerpoint"
            elif file_name.endswith(('.zip', '.rar', '.7zip')):
                return "archive"
            
            # Si no tiene extensión, intentar detectar por contenido
            try:
                # Verificar si es un PDF por el contenido
                if hasattr(self.file, 'path') and os.path.exists(self.file.path):
                    with open(self.file.path, 'rb') as f:
                        header = f.read(4)
                        if header == b'%PDF':
                            return "pdf"
            except:
                pass
        
        # Si no se pudo detectar por extensión o contenido, 
        # asumir que es PDF si no tiene extensión (para archivos antiguos)
        if self.file and hasattr(self.file, 'name') and self.file.name:
            if '.' not in self.file.name:
                # Archivo sin extensión, asumir PDF para compatibilidad
                return "pdf"
        
        return "file"
    
    def get_file_url(self):
        """Retorna la URL del archivo, ya sea local o externa"""
        if self.external_url:
            # Para Google Drive, convertir a formato de preview si es necesario
            if 'drive.google.com' in self.external_url:
                # Si es un enlace de visualización, convertirlo a preview
                if '/view?' in self.external_url:
                    return self.external_url.replace('/view?', '/preview?')
                # Si ya es preview, devolverlo tal como está
                elif '/preview?' in self.external_url:
                    return self.external_url
                # Si no tiene parámetros, agregar preview
                else:
                    return self.external_url + '?usp=preview'
            return self.external_url
        elif self.file and hasattr(self.file, 'url'):
            try:
                return self.file.url
            except ValueError:
                # El archivo no existe o no está asociado
                return None
        return None
    
    def is_completed_by(self, user):
        """Verifica si el documento ha sido completado por el usuario"""
        return DocumentCompletion.objects.filter(
            user=user, 
            document=self
        ).exists()
    
    def mark_as_completed(self, user):
        """Marca el documento como completado por el usuario"""
        DocumentCompletion.objects.get_or_create(
            user=user,
            document=self
        )
    
    def mark_as_incomplete(self, user):
        """Marca el documento como incompleto por el usuario"""
        DocumentCompletion.objects.filter(
            user=user,
            document=self
        ).delete()

    def has_embedded_videos(self):
        """
        Detecta si el PPTX tiene videos embebidos
        """
        if not self.file.name.lower().endswith(('.ppt', '.pptx')):
            return False
        
        try:
            # Abrir el archivo PPTX como ZIP
            with zipfile.ZipFile(self.file.path, 'r') as zip_file:
                # Buscar archivos de video en la estructura del PPTX
                video_extensions = ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.webm']
                
                for file_info in zip_file.filelist:
                    filename = file_info.filename.lower()
                    # Buscar en media/embeddings/ (donde PowerPoint guarda videos)
                    if 'media/embeddings/' in filename:
                        return True
                    # Buscar archivos con extensiones de video
                    if any(filename.endswith(ext) for ext in video_extensions):
                        return True
                
                return False
        except Exception:
            return False

    def get_slide_count(self):
        """
        Obtiene el número total de diapositivas
        """
        if not self.file.name.lower().endswith(('.ppt', '.pptx')):
            return 0
        
        try:
            prs = Presentation(self.file.path)
            return len(prs.slides)
        except Exception:
            return 0

    def convert_pptx_to_html(self):
        """
        Convierte PPTX sin videos a HTML para visualización mejorada
        """
        if not self.file.name.lower().endswith(('.ppt', '.pptx')):
            return None
        
        try:
            prs = Presentation(self.file.path)
            slides_html = []
            
            for i, slide in enumerate(prs.slides):
                slide_html = f"""
                <div class="slide" id="slide-{i+1}" style="display: {'block' if i == 0 else 'none'};">
                    <div class="slide-content">
                        <h3 class="slide-number">Diapositiva {i+1}</h3>
                        <div class="slide-elements">
                """
                
                # Procesar cada forma en la diapositiva
                for shape in slide.shapes:
                    slide_html += self._process_shape(shape)
                
                slide_html += """
                        </div>
                    </div>
                </div>
                """
                slides_html.append(slide_html)
            
            return slides_html
        except Exception as e:
            print(f"Error converting PPTX: {e}")
            return None

    def _process_shape(self, shape):
        """
        Procesa una forma individual y retorna su HTML
        """
        html = ""
        
        # Texto con formato
        if hasattr(shape, 'text_frame') and shape.text_frame:
            html += self._process_text_frame(shape.text_frame)
        
        # Imágenes
        elif hasattr(shape, 'image'):
            html += self._process_image(shape)
        
        # Tablas
        elif hasattr(shape, 'table'):
            html += self._process_table(shape.table)
        
        # Formas y gráficos
        elif hasattr(shape, 'shape_type'):
            html += self._process_graphic(shape)
        
        return html

    def _process_text_frame(self, text_frame):
        """
        Procesa texto con formato
        """
        html = ""
        
        for paragraph in text_frame.paragraphs:
            # Determinar el nivel de encabezado basado en el tamaño de fuente
            font_size = 18  # tamaño por defecto
            if paragraph.runs:
                font_size = paragraph.runs[0].font.size.pt if paragraph.runs[0].font.size else 18
            
            # Aplicar estilo de encabezado
            if font_size >= 24:
                tag = "h2"
            elif font_size >= 20:
                tag = "h3"
            elif font_size >= 16:
                tag = "h4"
            else:
                tag = "p"
            
            # Procesar el párrafo
            paragraph_html = f"<{tag} class='slide-text'>"
            
            for run in paragraph.runs:
                # Aplicar formato
                text = run.text
                if run.font.bold:
                    text = f"<strong>{text}</strong>"
                if run.font.italic:
                    text = f"<em>{text}</em>"
                if run.font.underline:
                    text = f"<u>{text}</u>"
                
                # Color de texto
                if run.font.color.rgb:
                    color = f"#{run.font.color.rgb:06x}"
                    text = f"<span style='color: {color}'>{text}</span>"
                
                paragraph_html += text
            
            paragraph_html += f"</{tag}>"
            html += paragraph_html
        
        return html

    def _process_image(self, shape):
        """
        Procesa imágenes de la diapositiva
        """
        try:
            img_stream = io.BytesIO()
            shape.image.save(img_stream, format='PNG')
            img_data = base64.b64encode(img_stream.getvalue()).decode()
            
            # Obtener dimensiones si están disponibles
            width = shape.width if hasattr(shape, 'width') else None
            height = shape.height if hasattr(shape, 'height') else None
            
            style = ""
            if width and height:
                # Convertir de EMU a píxeles (1 EMU = 1/914400 pulgada)
                width_px = int(width / 914400 * 96)  # 96 DPI
                height_px = int(height / 914400 * 96)
                style = f"style='max-width: {width_px}px; max-height: {height_px}px;'"
            
            return f'<img src="data:image/png;base64,{img_data}" class="slide-image" alt="Imagen de diapositiva" {style}>'
        except Exception as e:
            print(f"Error processing image: {e}")
            return ""

    def _process_table(self, table):
        """
        Procesa tablas de la diapositiva
        """
        html = '<table class="slide-table">'
        
        for row in table.rows:
            html += '<tr>'
            for cell in row.cells:
                html += f'<td class="slide-table-cell">{cell.text}</td>'
            html += '</tr>'
        
        html += '</table>'
        return html

    def _process_graphic(self, shape):
        """
        Procesa formas y gráficos
        """
        # Para formas básicas, crear un contenedor
        shape_types = {
            1: "rectángulo",
            2: "rectángulo redondeado", 
            3: "óvalo",
            4: "diamante",
            5: "triángulo",
            6: "flecha",
            7: "línea",
            8: "forma libre"
        }
        
        shape_name = shape_types.get(shape.shape_type, "forma")
        
        # Si tiene texto, procesarlo
        if hasattr(shape, 'text_frame') and shape.text_frame:
            return f'<div class="slide-shape {shape_name.lower().replace(" ", "-")}">{self._process_text_frame(shape.text_frame)}</div>'
        
        return f'<div class="slide-shape {shape_name.lower().replace(" ", "-")}"></div>'

    def delete(self, *args, **kwargs):
        self.file.delete(save=False)
        super().delete(*args, **kwargs)


@receiver(post_save, sender=Upload)
def log_upload_save(sender, instance, created, **kwargs):
    if created:
        message = _(
            f"The file '{instance.title}' has been uploaded to the course '{instance.course}'."
        )
    else:
        message = _(
            f"The file '{instance.title}' of the course '{instance.course}' has been updated."
        )
    ActivityLog.objects.create(message=message)


@receiver(post_delete, sender=Upload)
def log_upload_delete(sender, instance, **kwargs):
    ActivityLog.objects.create(
        message=_(
            f"The file '{instance.title}' of the course '{instance.course}' has been deleted."
        )
    )


# Modelo actualizado para manejar URLs de Vimeo
class UploadVideo(models.Model):
    title = models.CharField(max_length=100)
    slug = models.SlugField(unique=True, blank=True)
    course = models.ForeignKey(Course, on_delete=models.CASCADE)
    vimeo_url = models.URLField(
        max_length=500,
        help_text=_("Ingresa la URL del video de Vimeo"),
        null=True,
        blank=True,
    )
    youtube_url = models.URLField(
        max_length=500,
        help_text=_("Ingresa la URL de YouTube"),
        null=True,
        blank=True,
    )
    order = models.PositiveIntegerField(
        default=0,
        help_text=_("Orden en que aparecerá el video (0 = primero)"),
    )
    summary = models.TextField(blank=True)
    timestamp = models.DateTimeField(auto_now_add=True)
    completed_by = models.ManyToManyField(
        settings.AUTH_USER_MODEL,
        through='VideoCompletion',
        related_name='completed_videos',
        blank=True
    )

    # Opcional: Si deseas mantener el campo 'video' para los registros existentes
    video = models.FileField(
        upload_to="course_videos/",
        help_text=_("Valid video formats: mp4, mkv, wmv, 3gp, f4v, avi, mp3"),
        validators=[
            FileExtensionValidator(["mp4", "mkv", "wmv", "3gp", "f4v", "avi", "mp3"])
        ],
        null=True,
        blank=True,
    )

    class Meta:
        ordering = ['order', 'timestamp']

    def __str__(self):
        return f"{self.title}"

    def get_absolute_url(self):
        return reverse(
            "video_single", kwargs={"slug": self.course.slug, "video_slug": self.slug}
        )

    def get_vimeo_id(self):
        """Extrae el ID del video de Vimeo desde la URL"""
        if not self.vimeo_url:
            return None
        try:
            # Patrones para diferentes formatos de URL de Vimeo
            patterns = [
                r'vimeo\.com/(?:video/)?(\d+)',  # https://vimeo.com/123456 o https://vimeo.com/video/123456
                r'player\.vimeo\.com/video/(\d+)',  # https://player.vimeo.com/video/123456
            ]
            for pattern in patterns:
                match = re.search(pattern, self.vimeo_url)
                if match:
                    vimeo_id = match.group(1)
                    # Validar que sea un número
                    if vimeo_id.isdigit():
                        return vimeo_id
        except Exception:
            pass
        return None

    def get_youtube_id(self):
        """Extrae el ID del video de YouTube desde la URL, manejando parámetros adicionales"""
        if not self.youtube_url:
            return None
        try:
            # Limpiar la URL de espacios y caracteres extra
            url = self.youtube_url.strip()
            
            # Patrones comunes de URLs de YouTube (más flexibles)
            patterns = [
                r'(?:youtube\.com/watch\?v=)([a-zA-Z0-9_-]{11})',  # https://www.youtube.com/watch?v=ABC123XYZ45
                r'(?:youtu\.be/)([a-zA-Z0-9_-]{11})',  # https://youtu.be/ABC123XYZ45
                r'(?:youtube\.com/embed/)([a-zA-Z0-9_-]{11})',  # https://www.youtube.com/embed/ABC123XYZ45
                r'(?:youtube\.com/v/)([a-zA-Z0-9_-]{11})',  # https://www.youtube.com/v/ABC123XYZ45
                r'(?:youtube\.com/watch\?.*v=)([a-zA-Z0-9_-]{11})',  # Con parámetros adicionales
            ]
            for pattern in patterns:
                match = re.search(pattern, url)
                if match:
                    youtube_id = match.group(1)
                    # Validar que tenga la longitud correcta (11 caracteres) y caracteres válidos
                    if len(youtube_id) == 11 and all(c.isalnum() or c in ['-', '_'] for c in youtube_id):
                        return youtube_id
        except Exception:
            pass
        return None
    
    def get_youtube_embed_url(self):
        """Genera la URL completa de embed de YouTube con parámetros necesarios"""
        youtube_id = self.get_youtube_id()
        if not youtube_id:
            return None
        # URL base con parámetros importantes para evitar Error 153
        # Usar el formato más compatible: /embed/VIDEO_ID
        base_url = f"https://www.youtube.com/embed/{youtube_id}"
        # Agregar parámetros para mejor compatibilidad y evitar errores
        # IMPORTANTE: Estos parámetros ayudan a evitar el Error 153
        params = [
            "rel=0",  # No mostrar videos relacionados al final
            "modestbranding=1",  # Menos branding de YouTube
            "playsinline=1",  # Reproducir inline en móviles
            "enablejsapi=1",  # Habilitar API de JavaScript (puede ayudar con algunos errores)
        ]
        return f"{base_url}?{'&'.join(params)}"

    def get_video_mime_type(self):
        """Detecta el tipo MIME del video basado en la extensión del archivo"""
        if not self.video:
            return "video/mp4"  # Por defecto
        
        try:
            filename = self.video.name.lower()
            # Mapeo de extensiones a tipos MIME
            mime_types = {
                '.mp4': 'video/mp4',
                '.mkv': 'video/x-matroska',
                '.wmv': 'video/x-ms-wmv',
                '.3gp': 'video/3gpp',
                '.f4v': 'video/x-f4v',
                '.avi': 'video/x-msvideo',
                '.webm': 'video/webm',
                '.ogg': 'video/ogg',
                '.mov': 'video/quicktime',
                '.flv': 'video/x-flv',
            }
            
            # Buscar la extensión en el nombre del archivo
            for ext, mime_type in mime_types.items():
                if filename.endswith(ext):
                    return mime_type
            
            # Si no se encuentra, intentar detectar por extensión genérica
            if '.' in filename:
                ext = '.' + filename.split('.')[-1]
                # Tipo genérico para videos desconocidos
                return f'video/{ext[1:]}' if ext[1:] else 'video/mp4'
        except Exception:
            pass
        
        return "video/mp4"  # Por defecto si hay algún error

    def is_completed_by(self, user):
        return self.completed_by.filter(id=user.id).exists()
    
    def mark_as_incomplete(self, user):
        """Marca el video como incompleto por el usuario"""
        VideoCompletion.objects.filter(
            user=user,
            video=self
        ).delete()

    def save(self, *args, **kwargs):
        if not self.slug:
            self.slug = unique_slug_generator(self)
        super().save(*args, **kwargs)

    def delete(self, *args, **kwargs):
        # Si tienes archivos de video almacenados, elimínalos al eliminar el registro
        if self.video:
            self.video.delete(save=False)
        super().delete(*args, **kwargs)


# Modelo para rastrear la finalización de documentos
class DocumentCompletion(models.Model):
    user = models.ForeignKey(
        settings.AUTH_USER_MODEL,
        on_delete=models.CASCADE,
        verbose_name=_("Usuario")
    )
    document = models.ForeignKey(
        Upload,
        on_delete=models.CASCADE,
        verbose_name=_("Documento")
    )
    completed_at = models.DateTimeField(
        auto_now_add=True,
        verbose_name=_("Completado el")
    )
    
    class Meta:
        unique_together = ['user', 'document']
        verbose_name = _("Finalización de Documento")
        verbose_name_plural = _("Finalizaciones de Documentos")
    
    def __str__(self):
        return f"{self.user.username} - {self.document.title}"


@receiver(pre_save, sender=UploadVideo)
def video_pre_save_receiver(sender, instance, **kwargs):
    if not instance.slug:
        instance.slug = unique_slug_generator(instance)


@receiver(post_save, sender=UploadVideo)
def log_uploadvideo_save(sender, instance, created, **kwargs):
    if created:
        message = _(
            f"El video '{instance.title}' ha sido agregado al curso '{instance.course}'."
        )
    else:
        message = _(
            f"El video '{instance.title}' del curso '{instance.course}' ha sido actualizado."
        )
    ActivityLog.objects.create(message=message)


@receiver(post_delete, sender=UploadVideo)
def log_uploadvideo_delete(sender, instance, **kwargs):
    ActivityLog.objects.create(
        message=_(
            f"El video '{instance.title}' del curso '{instance.course}' ha sido eliminado."
        )
    )


class VideoCompletion(models.Model):
    user = models.ForeignKey(settings.AUTH_USER_MODEL, on_delete=models.CASCADE)
    video = models.ForeignKey(UploadVideo, on_delete=models.CASCADE)
    completed_at = models.DateTimeField(auto_now_add=True)

    class Meta:
        unique_together = ('user', 'video')
        ordering = ['-completed_at']

    def __str__(self):
        return f"{self.user.username} completed {self.video.title}"


class CourseOffer(models.Model):
    """NOTA: Solo el jefe de departamento puede ofrecer cursos semestrales"""

    dep_head = models.ForeignKey("accounts.DepartmentHead", on_delete=models.CASCADE)

    def __str__(self):
        return str(self.dep_head)
